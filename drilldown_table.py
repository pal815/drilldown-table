"""
drilldown_table.py — 드릴다운(corner-header / ┌, 일명 ㄱ자) 계층형 표 생성기
================================================================================
공용 IR(Grid) 한 번 계산 → 백엔드 3개(Excel/Word/PowerPoint)가 동일 양식으로 렌더.

양식(v4.1, 사용자 확정)
- 라벨 2칸 분리(항목명 띠 + 세부 들여쓰기). 항목명 = HEAD_FILL 세로 띠로 블록 관통(ㄱ 세로획).
- 헤드라인(항목명+합계) = 가로 띠(ㄱ 가로획). '계' 라벨 없음.
- 세부(SUB_FILL) = medium 안쪽 ┌로 한 번 더 감싼 부분집합. navy = 코너 1칸.
- drill-through(Excel 한정) = 아웃라인 그룹 ±버튼.

orientation: "row" / "column" / "both"(행·열 모두 ㄱ자, 계 자동합산)

API
    build_xlsx(model, path, orientation)          # 단일 시트
    build_workbook(variants, path)                # 다중 시트 [(model, orient, title?)]
    build_docx(variants, path)                    # 표들을 한 문서에
    build_pptx(variants, path)                    # 변형당 1슬라이드
    layout(model, orientation) -> Grid            # IR만 필요할 때
"""
from __future__ import annotations

import json
import re
import sys

# ── 팔레트 / 타이포 ───────────────────────────────────────────────────────────
FONT = "맑은 고딕"
C_NAVY = "1F3864"
HEAD_FILL = "D6E0F2"
SUB_FILL = "EDF2FA"
C_WHITE = "FFFFFF"
C_GRID = "C9D2E0"
C_HEAVY = "404040"
C_TEXT = "1A1A1A"
C_MUTE = "44546A"

THIN, HEAVY = "thin", "medium"          # 테두리 굵기(문자열 IR)
_PCT = re.compile(r"^-?\d+(?:\.(\d+))?%$")


class _Al:
    __slots__ = ("h", "ind", "wrap")

    def __init__(self, h, ind=0, wrap=False):
        self.h, self.ind, self.wrap = h, ind, wrap


A_CENTER = _Al("center", wrap=True)
A_LEFT = _Al("left", wrap=True)
A_RIGHT = _Al("right")
A_LEFT_IND = _Al("left", 1)

# ── 테마(팔레트) ──────────────────────────────────────────────────────────────
THEME = {
    "color": dict(navy=C_NAVY, head=HEAD_FILL, sub=SUB_FILL, white=C_WHITE, text=C_TEXT, mute=C_MUTE, grid=C_GRID, heavy=C_HEAVY),
    "grey":  dict(navy="595959", head="D9D9D9", sub="F2F2F2", white=C_WHITE, text="1A1A1A", mute="404040", grid="BFBFBF", heavy="262626"),
    "mono":  dict(navy=None, head=None, sub=None, white="000000", text="000000", mute="333333", grid="808080", heavy="000000"),
}


def _interp(a, b, t):
    a, b = a.lstrip("#"), b.lstrip("#")
    return "".join("%02X" % round(int(a[i:i + 2], 16) + (int(b[i:i + 2], 16) - int(a[i:i + 2], 16)) * t)
                   for i in (0, 2, 4))


def _level_fill(k, D):
    """단계 k(0=최상위) 채움 — HEAD_FILL→SUB_FILL 보간. depth=2면 정확히 HEAD/SUB(회귀 0)."""
    if D <= 1:
        return HEAD_FILL
    return _interp(HEAD_FILL, SUB_FILL, min(k, D - 1) / (D - 1))


def themed(grid, theme):
    """color 테마 hex를 다른 테마로 변환(fill·color). 단계 채움 램프도 테마별로 재매핑. 'color'면 그대로."""
    if theme == "color":
        return grid
    T = THEME[theme]
    D = max(getattr(grid, "_depth", 2), 2)
    fmap = {C_NAVY: T["navy"]}
    for k in range(D):                                   # 단계별 채움 램프(2단계면 HEAD/SUB만)
        fmap[_level_fill(k, D)] = None if T["head"] is None else _interp(T["head"], T["sub"], min(k, D - 1) / (D - 1))
    cmap = {C_WHITE: T["white"], C_TEXT: T["text"], C_MUTE: T["mute"]}
    for cell in grid.g.values():
        if cell.fill in fmap:
            cell.fill = fmap[cell.fill]
        if cell.color in cmap:
            cell.color = cmap[cell.color]
    return grid


# ── IR ────────────────────────────────────────────────────────────────────────
class GCell:
    __slots__ = ("value", "text", "numfmt", "fill", "bold", "size",
                 "halign", "indent", "wrap", "color", "top", "left", "right", "bottom")

    def __init__(self):
        self.value = None
        self.text = ""
        self.numfmt = None
        self.fill = None
        self.bold = False
        self.size = 10
        self.halign = None
        self.indent = 0
        self.wrap = False
        self.color = C_TEXT
        self.top = self.left = self.right = self.bottom = None


class Grid:
    def __init__(self):
        self.g = {}
        self.merges = []                 # 실제 병합(전 백엔드)
        self.soft_merges = []            # 시각 병합(docx/pptx 텍스트 spanning 전용, xlsx 무시)
        self.widths = {}
        self.heights = {}
        self.outline_rows = {}
        self.outline_cols = {}
        self.summary_below = None
        self.summary_right = None
        self.freeze = None
        self.nr = self.nc = 0

    def cell(self, r, c):
        k = (r, c)
        cell = self.g.get(k)
        if cell is None:
            cell = self.g[k] = GCell()
        self.nr = max(self.nr, r)
        self.nc = max(self.nc, c)
        return cell


# ── 저수준 헬퍼(IR 대상) ──────────────────────────────────────────────────────
def _looks_num(v) -> bool:
    s = str(v).strip().replace(",", "").rstrip("%")
    try:
        float(s)
        return True
    except ValueError:
        return False


def _disp(value):
    """입력값 → (xlsx_value, numfmt, display_text, halign)."""
    if isinstance(value, str):
        mm = _PCT.match(value.strip())
        if mm:
            f = float(value.strip()[:-1]) / 100.0
            dec = len(mm.group(1)) if mm.group(1) else 0
            return f, ("0%" if dec == 0 else "0." + "0" * dec + "%"), value.strip(), "right"
        return value, None, value, ("right" if _looks_num(value) else "left")
    if isinstance(value, bool):
        return value, None, str(value), "left"
    if isinstance(value, (int, float)):
        return value, "#,##0", f"{value:,}", "right"
    return value, None, ("" if value is None else str(value)), None


def _fmt(cell, value):
    cell.value, cell.numfmt, cell.text, ha = _disp(value)
    if ha:
        cell.halign = ha


def _lbl(grid, r, c, label):
    cell = grid.cell(r, c)
    if label is not None and label != "":
        cell.value = label
        cell.text = str(label)
    return cell


def _style(cell, *, bold=False, size=10, color=C_TEXT, fill=None, align=None):
    cell.bold, cell.size, cell.color = bold, size, color
    if fill:
        cell.fill = fill
    if align:
        cell.halign, cell.indent, cell.wrap = align.h, align.ind, align.wrap


def _merge(grid, r1, c1, r2, c2):
    grid.merges.append((r1, c1, r2, c2))
    grid.cell(r2, c2)


def _grid_borders(grid, r1, c1, r2, c2):
    for r in range(r1, r2 + 1):
        for c in range(c1, c2 + 1):
            cell = grid.cell(r, c)
            cell.top = cell.left = cell.right = cell.bottom = THIN


def _edges(cell, top=None, left=None, right=None, bottom=None):
    cell.top = top or cell.top
    cell.left = left or cell.left
    cell.right = right or cell.right
    cell.bottom = bottom or cell.bottom


def _box(cell, t=None, l=None, r=None, b=None):
    cell.top, cell.left, cell.right, cell.bottom = t, l, r, b


def _clear_left(cell):
    cell.left = None


def _titles(grid, m, r, ncol):
    if m.get("title"):
        _merge(grid, r, 1, r, max(ncol, 2))
        _style(_lbl(grid, r, 1, m["title"]), bold=True, size=13, align=A_LEFT)
        grid.heights[r] = 22
        r += 1
    if m.get("unit"):
        _merge(grid, r, 1, r, max(ncol, 2))
        _style(_lbl(grid, r, 1, m["unit"]), size=9, color=C_MUTE, align=A_RIGHT)
        r += 1
    return r


def _outer(grid, r1, c1, r2, c2):
    for c in range(c1, c2 + 1):
        _edges(grid.cell(r1, c), top=HEAVY)
        _edges(grid.cell(r2, c), bottom=HEAVY)
    for r in range(r1, r2 + 1):
        _edges(grid.cell(r, c1), left=HEAVY)
        _edges(grid.cell(r, c2), right=HEAVY)


def _vband(grid, col, b0, b1):
    _box(grid.cell(b0, col), t=HEAVY, l=HEAVY)
    for rr in range(b0 + 1, b1 + 1):
        _box(grid.cell(rr, col), l=HEAVY, r=HEAVY, b=(HEAVY if rr == b1 else None))


def _hband(grid, row, c0, c1):       # _vband 의 전치(가로 띠)
    _box(grid.cell(row, c0), t=HEAVY, l=HEAVY)
    for cc in range(c0 + 1, c1 + 1):
        _box(grid.cell(row, cc), t=HEAVY, b=HEAVY, r=(HEAVY if cc == c1 else None))


# ── 다단계 공용 헬퍼 ──────────────────────────────────────────────────────────
def _normalize_items(m):
    """기존 2단계 items → 재귀 nodes 트리(1:1). depth=2를 N=2 특수케이스로 흡수."""
    out = []
    for it in m["items"]:
        out.append({
            "label": it["label"],
            "summary": it.get("summary"),
            "children": [{"label": d["label"], "values": d["values"], "children": None, "summary": None}
                         for d in it.get("details", [])],
        })
    return out


def _nodes_of(m):
    return m.get("nodes") or _normalize_items(m)


def _max_depth(nodes):
    return max((1 if not n.get("children") else 1 + _max_depth(n["children"])) for n in nodes)


def _tree_size(nodes):
    """전체 노드 수(열 기준에서 각 노드 = 1열)."""
    return sum(1 + (_tree_size(n["children"]) if n.get("children") else 0) for n in nodes)


# ── 레이아웃: 행 기준 (N단계) ─────────────────────────────────────────────────
def _build_row(grid, m):
    nodes = _nodes_of(m)
    attrs = m["attributes"]
    A = len(attrs)
    DEPTH = _max_depth(nodes)
    grid._depth = DEPTH                      # themed() 단계 램프용
    levels = m.get("level_labels")          # 다단계: 단계명 리스트(없으면 코너 1개=2단계 호환)
    C0, C_AT0, ncol = 1, 1 + DEPTH, DEPTH + A
    r = _titles(grid, m, 1, ncol)

    hdr = r
    if levels:
        for d in range(DEPTH):
            _style(_lbl(grid, hdr, C0 + d, levels[d] if d < len(levels) else ""),
                   bold=True, color=C_WHITE, fill=C_NAVY, align=A_LEFT)
    else:
        _style(_lbl(grid, hdr, C0, m.get("corner_label", "")), bold=True, color=C_WHITE, fill=C_NAVY, align=A_LEFT)
        for d in range(1, DEPTH):
            _style(grid.cell(hdr, C0 + d), fill=C_NAVY)
    for j, a in enumerate(attrs):
        _style(_lbl(grid, hdr, C_AT0 + j, a), bold=True, color=C_WHITE, fill=C_NAVY, align=A_CENTER)

    st = {"r": hdr + 1}
    flat = []                                # (level, r0, r1, is_group)

    def emit(node, level):
        lcol = C0 + level
        is_grp = bool(node.get("children"))
        is_hdr = level < DEPTH - 1                          # 최하위 단계가 아니면(그룹/잎 무관) 헤더 = bold
        r0 = st["r"]
        sfill = _level_fill(level, DEPTH)
        lcolor = C_TEXT if level == 0 else C_MUTE
        for cc in range(C0, lcol):                          # 왼쪽 조상 띠칸 = 조상 단계 색
            _style(grid.cell(r0, cc), fill=_level_fill(cc - C0, DEPTH))
        _style(_lbl(grid, r0, lcol, node["label"]), bold=is_hdr, fill=sfill, color=lcolor,
               align=(A_LEFT if level == 0 else _Al("left", level)))
        for cc in range(lcol + 1, C_AT0):                  # 라벨칸 우측(헤드라인/연속) = 자기 단계 색
            _style(grid.cell(r0, cc), bold=is_hdr, fill=sfill)
        vals = node.get("summary") if is_grp else node.get("values")
        for j in range(A):
            if vals is not None and vals[j] is not None:
                _fmt(grid.cell(r0, C_AT0 + j), vals[j])
            _style(grid.cell(r0, C_AT0 + j), bold=is_hdr, fill=sfill)
        st["r"] += 1
        if is_grp:
            for ch in node["children"]:
                emit(ch, level + 1)
        flat.append((level, r0, st["r"] - 1, is_grp))

    for n in nodes:
        emit(n, 0)
    last = st["r"] - 1

    _grid_borders(grid, hdr, C0, last, ncol)
    _outer(grid, hdr, C0, last, ncol)
    if not levels:                           # 2단계 호환: 코너 칸막이 투명
        grid.cell(hdr, C0).right = None
        grid.cell(hdr, C0 + 1).left = None
    for (level, r0, r1, is_grp) in flat:     # Pass A: 헤더 단계 헤드라인 상단 + 하단 구분/블록선
        if level < DEPTH - 1:
            lcol = C0 + level
            for c in range(lcol, ncol + 1):
                _edges(grid.cell(r0, c), top=HEAVY)
            lo = lcol + 1 if is_grp else lcol            # 그룹=summary↔children / 헤더잎=단일행 블록 하단
            for c in range(lo, ncol + 1):
                _edges(grid.cell(r0, c), bottom=HEAVY)
    for (level, r0, r1, is_grp) in flat:     # Pass B: 세로 띠 / 최하위 잎 좌측 + 아웃라인
        lcol = C0 + level
        if level < DEPTH - 1:
            if is_grp:
                _vband(grid, lcol, r0, r1)               # 다행 그룹 띠(_box 덮어쓰므로 Pass A 뒤)
            else:
                _edges(grid.cell(r0, lcol), left=HEAVY)  # 단일행 헤더잎: _box 회피로 top/bottom 보존
            _clear_left(grid.cell(r0, lcol + 1))
        else:
            _edges(grid.cell(r0, lcol), left=HEAVY)
        if level >= 1:
            grid.outline_rows[r0] = level
    grid.summary_below = False

    grid.widths[C0] = 12
    for d in range(1, DEPTH):
        grid.widths[C0 + d] = 14 if d == DEPTH - 1 else 12
    for j in range(A):
        grid.widths[C_AT0 + j] = 12
    grid.freeze = (hdr + 1, C_AT0)


# ── 레이아웃: 열 기준 ─────────────────────────────────────────────────────────
def _build_column(grid, m):
    if "nodes" in m:                         # 다단계(3+) → 전치 빌더
        return _build_column_ndepth(grid, m)
    return _build_column_2level(grid, m)      # 기존 2단계(items) — 회귀 0


def _build_column_ndepth(grid, m):
    """열 기준 N단계 = 행 기준 N단계의 전치. 단계=라벨 행, 항목=열, 값=아래로."""
    nodes = m["nodes"]
    attrs = m["attributes"]
    A = len(attrs)
    DEPTH = _max_depth(nodes)
    grid._depth = DEPTH
    levels = m.get("level_labels")
    C_ALBL = 1
    ncol = 1 + _tree_size(nodes)
    r = _titles(grid, m, 1, ncol)
    R0 = r                                    # 첫 단계 라벨 행
    R_DATA0 = R0 + DEPTH
    last_r = R_DATA0 + A - 1

    if levels:                                # 코너(네이비): 단계명 세로 + 속성명
        for k in range(DEPTH):
            _style(_lbl(grid, R0 + k, C_ALBL, levels[k] if k < len(levels) else ""),
                   bold=True, color=C_WHITE, fill=C_NAVY, align=A_LEFT_IND)
    else:
        _merge(grid, R0, C_ALBL, R0 + DEPTH - 1, C_ALBL)
        _style(_lbl(grid, R0, C_ALBL, m.get("corner_label", "")), bold=True, color=C_WHITE, fill=C_NAVY, align=A_CENTER)
    for i, a in enumerate(attrs):
        _style(_lbl(grid, R_DATA0 + i, C_ALBL, a), bold=True, color=C_WHITE, fill=C_NAVY, align=A_LEFT_IND)

    st = {"c": 2}
    flat = []

    def emit(node, level):
        is_grp = bool(node.get("children"))
        is_hdr = level < DEPTH - 1
        c0 = st["c"]
        sfill = _level_fill(level, DEPTH)
        lcolor = C_TEXT if level == 0 else C_MUTE
        for rr in range(R0, R0 + level):                     # 위 조상 라벨행 = 조상 단계 색
            _style(grid.cell(rr, c0), fill=_level_fill(rr - R0, DEPTH))
        _style(_lbl(grid, R0 + level, c0, node["label"]), bold=is_hdr, fill=sfill, color=lcolor,
               align=(A_LEFT if level == 0 else A_CENTER))
        for rr in range(R0 + level + 1, R0 + DEPTH):         # 아래 헤드라인 연속 = 자기 단계 색
            _style(grid.cell(rr, c0), bold=is_hdr, fill=sfill)
        vals = node.get("summary") if is_grp else node.get("values")
        for i in range(A):
            if vals is not None and vals[i] is not None:
                _fmt(grid.cell(R_DATA0 + i, c0), vals[i])
            _style(grid.cell(R_DATA0 + i, c0), bold=is_hdr, fill=sfill)
        st["c"] += 1
        if is_grp:
            for ch in node["children"]:
                emit(ch, level + 1)
        c1 = st["c"] - 1
        if c1 > c0:                                          # 그룹 라벨 가로 span
            grid.soft_merges.append((R0 + level, c0, R0 + level, c1))
            for cc in range(c0 + 1, c1 + 1):
                _style(grid.cell(R0 + level, cc), bold=is_hdr, fill=sfill)
        flat.append((level, c0, c1, is_grp))

    for n in nodes:
        emit(n, 0)
    last_c = st["c"] - 1

    _grid_borders(grid, R0, C_ALBL, last_r, last_c)
    _outer(grid, R0, C_ALBL, last_r, last_c)
    for cc in range(C_ALBL, last_c + 1):                     # 라벨↔데이터 구분선
        _edges(grid.cell(R0 + DEPTH - 1, cc), bottom=HEAVY)
    for (level, c0, c1, is_grp) in flat:                    # Pass A: 헤더 좌측선 + 우측 구분/블록(전치)
        if level < DEPTH - 1:
            lrow = R0 + level
            for rr in range(lrow, last_r + 1):
                _edges(grid.cell(rr, c0), left=HEAVY)
            lo = lrow + 1 if is_grp else lrow
            for rr in range(lo, last_r + 1):
                _edges(grid.cell(rr, c0), right=HEAVY)
    for (level, c0, c1, is_grp) in flat:                    # Pass B: 가로 띠 / 최하위 잎 상단
        lrow = R0 + level
        if level < DEPTH - 1:
            if is_grp:
                _hband(grid, lrow, c0, c1)
            else:
                _edges(grid.cell(lrow, c0), top=HEAVY)
            grid.cell(lrow + 1, c0).top = None
        else:
            _edges(grid.cell(lrow, c0), top=HEAVY)
        if level >= 1:
            grid.outline_cols[c0] = level
    grid.summary_right = False

    grid.widths[C_ALBL] = 13
    for cc in range(2, last_c + 1):
        grid.widths[cc] = 9
    grid.freeze = (R_DATA0, 2)


def _build_column_2level(grid, m):
    attrs = m["attributes"]
    A = len(attrs)
    C_ALBL = 1
    r = _titles(grid, m, 1, 1 + sum((1 if it.get("summary") is not None else 0) + len(it.get("details", []))
                                    for it in m["items"]))
    r_item, r_sub, r_data0 = r, r + 1, r + 2

    _merge(grid, r_item, C_ALBL, r_sub, C_ALBL)
    _style(_lbl(grid, r_item, C_ALBL, m.get("corner_label", "")), bold=True, color=C_WHITE, fill=C_NAVY, align=A_CENTER)

    c = 2
    blocks = []
    for it in m["items"]:
        c0 = c
        summ = it.get("summary")
        if summ is not None:
            _style(grid.cell(r_sub, c), fill=HEAD_FILL)
            for i, v in enumerate(summ):
                _fmt(grid.cell(r_data0 + i, c), v)
                _style(grid.cell(r_data0 + i, c), bold=True, fill=HEAD_FILL)
            c += 1
        sub_first = c
        for d in it.get("details", []):
            _style(_lbl(grid, r_sub, c, d["label"]), fill=SUB_FILL, color=C_MUTE, align=A_CENTER)
            for i, v in enumerate(d["values"]):
                _fmt(grid.cell(r_data0 + i, c), v)
                _style(grid.cell(r_data0 + i, c), fill=SUB_FILL)
            c += 1
        c1 = c - 1
        for cc in range(c0, c1 + 1):
            _style(_lbl(grid, r_item, cc, it["label"] if cc == c0 else None),
                   bold=True, fill=HEAD_FILL, align=(A_LEFT if cc == c0 else None))
        if c1 > c0:
            grid.soft_merges.append((r_item, c0, r_item, c1))
        blocks.append((c0, sub_first if summ is not None else c0, c1))

    last_c = c - 1
    for i, a in enumerate(attrs):
        _style(_lbl(grid, r_data0 + i, C_ALBL, a), bold=True, color=C_WHITE, fill=C_NAVY, align=A_LEFT_IND)
    last_r = r_data0 + A - 1

    _grid_borders(grid, r_item, C_ALBL, last_r, last_c)
    _outer(grid, r_item, C_ALBL, last_r, last_c)
    for cc in range(C_ALBL, last_c + 1):
        _edges(grid.cell(r_sub, cc), bottom=HEAVY)         # 헤더 하단선(계 포함 전 열): 서브라벨 ↔ 데이터 구분
    for c0, sub_first, c1 in blocks:
        for cc in range(c0, c1 + 1):
            _box(grid.cell(r_item, cc), t=HEAVY, l=(HEAVY if cc == c0 else None), r=(HEAVY if cc == c1 else None))
        for rr in range(r_item, last_r + 1):
            _edges(grid.cell(rr, c0), left=HEAVY)
            _edges(grid.cell(rr, c1), right=HEAVY)
        if sub_first > c0:
            for rr in range(r_sub, last_r + 1):
                _edges(grid.cell(rr, sub_first), left=HEAVY)
            for cc in range(sub_first, c1 + 1):
                _edges(grid.cell(r_sub, cc), top=HEAVY)
            grid.cell(r_sub, c0).top = None   # 항목명 ↔ 계 서브라벨만 연속(아래 헤더선 HEAVY는 유지)
    for cc in range(2, last_c + 1):
        if not any(cc == c0 for c0, _sf, _c1 in blocks):
            grid.outline_cols[cc] = 1
    grid.summary_right = False

    grid.widths[C_ALBL] = 13
    for cc in range(2, last_c + 1):
        grid.widths[cc] = 9
    grid.freeze = (r_data0, 2)


# ── 레이아웃: 행 + 열 ─────────────────────────────────────────────────────────
def _build_both(grid, m):
    rg, cg, data = m["row_groups"], m["col_groups"], m["data"]
    cg_det = {c["label"]: c["details"] for c in cg}
    C_GRP, C_DET = 1, 2
    r = _titles(grid, m, 1, 2 + sum(1 + len(c["details"]) for c in cg))
    r_citem, r_csub, r0 = r, r + 1, r + 2

    ccols, col_blocks = [], []
    c = 3
    for cgrp in cg:
        c0 = c
        ccols.append((c, cgrp["label"], None)); c += 1
        for cd in cgrp["details"]:
            ccols.append((c, cgrp["label"], cd)); c += 1
        col_blocks.append((c0, c0 + 1, c - 1))
    last_c = c - 1

    _merge(grid, r_citem, C_GRP, r_csub, C_DET)
    _style(_lbl(grid, r_citem, C_GRP, m.get("corner_label", "")), bold=True, color=C_WHITE, fill=C_NAVY, align=A_CENTER)
    for (cgrp, (c0, _sf, c1)) in zip(cg, col_blocks):
        for cc in range(c0, c1 + 1):                                  # 연도 행 전체 HEAD_FILL → 헤더 느낌(계 서브라벨과 동색)
            _style(_lbl(grid, r_citem, cc, cgrp["label"] if cc == c0 else None), bold=True, fill=HEAD_FILL, align=A_LEFT)
    for cidx, ci, cd in ccols:
        if cd is None:
            _style(grid.cell(r_csub, cidx), fill=HEAD_FILL)
        else:
            _style(_lbl(grid, r_csub, cidx, cd), fill=SUB_FILL, color=C_MUTE, align=A_CENTER)

    def v_fine(ri, rd, ci, cd):
        return data[ri][rd][ci][cd]

    def v_rd(ri, rd, ci, cd):
        return sum(v_fine(ri, rd, ci, x) for x in cg_det[ci]) if cd is None else v_fine(ri, rd, ci, cd)

    def v_head(rgrp, ci, cd):
        return sum(v_rd(rgrp["label"], rd, ci, cd) for rd in rgrp["details"])

    r = r0
    row_blocks = []
    for rgrp in rg:
        b0 = r
        _style(_lbl(grid, r, C_GRP, rgrp["label"]), bold=True, fill=HEAD_FILL, align=A_LEFT)
        _style(grid.cell(r, C_DET), bold=True, fill=HEAD_FILL)
        for cidx, ci, cd in ccols:
            _fmt(grid.cell(r, cidx), v_head(rgrp, ci, cd))
            _style(grid.cell(r, cidx), bold=True, fill=HEAD_FILL)
        r += 1
        det_rows = []
        for rd in rgrp["details"]:
            _style(grid.cell(r, C_GRP), fill=HEAD_FILL)
            _style(_lbl(grid, r, C_DET, rd), fill=SUB_FILL, color=C_MUTE, align=A_LEFT_IND)
            for cidx, ci, cd in ccols:
                _fmt(grid.cell(r, cidx), v_rd(rgrp["label"], rd, ci, cd))
                _style(grid.cell(r, cidx), fill=SUB_FILL)
            det_rows.append(r)
            r += 1
        row_blocks.append((b0, r - 1, det_rows))
    last_r = r - 1

    _grid_borders(grid, r_citem, C_GRP, last_r, last_c)
    _outer(grid, r_citem, C_GRP, last_r, last_c)
    for cc in range(C_GRP, last_c + 1):
        _edges(grid.cell(r_csub, cc), bottom=HEAVY)
    for c0, sub_first, c1 in col_blocks:
        for cc in range(c0, c1 + 1):
            _box(grid.cell(r_citem, cc), t=HEAVY, l=(HEAVY if cc == c0 else None), r=(HEAVY if cc == c1 else None))
        for rr in range(r_citem, last_r + 1):
            _edges(grid.cell(rr, c0), left=HEAVY)
        for rr in range(r_csub, last_r + 1):                    # 계↔세부 세로(서브라벨 행부터 — 연도 행 제외)
            _edges(grid.cell(rr, sub_first), left=HEAVY)
        for cc in range(sub_first, c1 + 1):
            _edges(grid.cell(r_csub, cc), top=HEAVY)            # 세부 ┌ 상단 가로
        grid.cell(r_csub, c0).top = None                        # 연도 ↔ 계 서브라벨 가로 투명(계 연속)
    for b0, b1, det_rows in row_blocks:
        for cc in range(C_GRP, last_c + 1):
            _edges(grid.cell(b0, cc), top=HEAVY)
        if det_rows:
            for cc in range(C_DET, last_c + 1):
                _edges(grid.cell(b0, cc), bottom=HEAVY)
        _vband(grid, C_GRP, b0, b1)
        _clear_left(grid.cell(b0, C_DET))
        for rr in det_rows:
            _edges(grid.cell(rr, C_DET), left=HEAVY)
            grid.outline_rows[rr] = 1
    for cidx, ci, cd in ccols:
        if cd is not None:
            grid.outline_cols[cidx] = 1
    grid.summary_below = grid.summary_right = False

    grid.widths[C_GRP] = 12
    grid.widths[C_DET] = 13
    for cc in range(3, last_c + 1):
        grid.widths[cc] = 8.5
    grid.freeze = (r0, 3)


_BUILDERS = {"row": _build_row, "column": _build_column, "both": _build_both}
_DEFTITLE = {"row": "행기준 ㄱ자", "column": "열기준 ㄱ자", "both": "행+열 ㄱ자"}


def layout(model, orientation="row") -> Grid:
    if orientation not in _BUILDERS:
        raise ValueError("orientation must be 'row', 'column' or 'both'")
    grid = Grid()
    _BUILDERS[orientation](grid, model)
    return grid


# ── 렌더러 공통: 병합 처리 ────────────────────────────────────────────────────
def _merge_maps(merges):
    """anchor (r,c) -> (r2,c2) 와 covered(비-anchor) 집합."""
    span = {}
    covered = set()
    for (r1, c1, r2, c2) in merges:
        span[(r1, c1)] = (r2, c2)
        for r in range(r1, r2 + 1):
            for c in range(c1, c2 + 1):
                if (r, c) != (r1, c1):
                    covered.add((r, c))
    return span, covered


_WEIGHT = {None: 0, THIN: 1, HEAVY: 2}


def _heaviest(a, b):
    return a if _WEIGHT[a] >= _WEIGHT[b] else b


def reconcile(grid):
    """공유 모서리의 양쪽 셀 테두리를 '굵은 쪽'으로 일치 (PPT/Word가 얇은 쪽 고르는 문제 해결)."""
    rec = {}
    for (r, c), cell in grid.g.items():
        top, left, right, bottom = cell.top, cell.left, cell.right, cell.bottom
        up = grid.g.get((r - 1, c))
        if up:
            top = _heaviest(top, up.bottom)
        lf = grid.g.get((r, c - 1))
        if lf:
            left = _heaviest(left, lf.right)
        rt = grid.g.get((r, c + 1))
        if rt:
            right = _heaviest(right, rt.left)
        dn = grid.g.get((r + 1, c))
        if dn:
            bottom = _heaviest(bottom, dn.top)
        rec[(r, c)] = (top, left, right, bottom)
    return rec


def _region_border(rec, r1, c1, r2, c2):
    tl = rec.get((r1, c1), (None, None, None, None))
    tr = rec.get((r1, c2), (None, None, None, None))
    bl = rec.get((r2, c1), (None, None, None, None))
    return (tl[0], tl[1], tr[2], bl[3])


def _is_dark(hexcolor):
    if not hexcolor:
        return False
    r, g, b = int(hexcolor[0:2], 16), int(hexcolor[2:4], 16), int(hexcolor[4:6], 16)
    return 0.299 * r + 0.587 * g + 0.114 * b < 110


# ── 렌더러: Excel ─────────────────────────────────────────────────────────────
def render_xlsx(grid, ws, *, grid_color=C_GRID, heavy_color=C_HEAVY):
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
    from openpyxl.utils import get_column_letter
    side = {THIN: Side(style="thin", color=grid_color), HEAVY: Side(style="medium", color=heavy_color)}
    ws.sheet_view.showGridLines = False
    for (r, c), cell in grid.g.items():
        x = ws.cell(r, c)
        if cell.value is not None:
            x.value = cell.value
        elif cell.text:
            x.value = cell.text
        if cell.numfmt:
            x.number_format = cell.numfmt
        x.font = Font(name=FONT, bold=cell.bold, size=cell.size, color=cell.color)
        if cell.fill:
            x.fill = PatternFill("solid", fgColor=cell.fill)
        if cell.halign:
            x.alignment = Alignment(horizontal=cell.halign, vertical="center", indent=cell.indent, wrap_text=cell.wrap)
        x.border = Border(top=side.get(cell.top), left=side.get(cell.left),
                          right=side.get(cell.right), bottom=side.get(cell.bottom))
    for (r1, c1, r2, c2) in grid.merges:
        ws.merge_cells(start_row=r1, start_column=c1, end_row=r2, end_column=c2)
    for c, w in grid.widths.items():
        ws.column_dimensions[get_column_letter(c)].width = w
    for r, h in grid.heights.items():
        ws.row_dimensions[r].height = h
    for r, lvl in grid.outline_rows.items():
        ws.row_dimensions[r].outlineLevel = lvl
    for c, lvl in grid.outline_cols.items():
        ws.column_dimensions[get_column_letter(c)].outlineLevel = lvl
    if grid.summary_below is not None:
        ws.sheet_properties.outlinePr.summaryBelow = grid.summary_below
    if grid.summary_right is not None:
        ws.sheet_properties.outlinePr.summaryRight = grid.summary_right
    if grid.freeze:
        ws.freeze_panes = ws.cell(*grid.freeze)


def build_workbook(variants, path, theme="color"):
    from openpyxl import Workbook
    T = THEME[theme]
    wb = Workbook()
    wb.remove(wb.active)
    for v in variants:
        model, orient = v[0], v[1]
        title = v[2] if len(v) > 2 else _DEFTITLE[orient]
        ws = wb.create_sheet(title)
        render_xlsx(themed(layout(model, orient), theme), ws, grid_color=T["grid"], heavy_color=T["heavy"])
    wb.save(path)
    return path


def build_xlsx(model, path, orientation="row", theme="color"):
    return build_workbook([(model, orientation)], path, theme)


def add_sheet(wb, model, orientation="row", title=None, theme="color"):
    T = THEME[theme]
    ws = wb.create_sheet(title or _DEFTITLE[orientation])
    render_xlsx(themed(layout(model, orientation), theme), ws, grid_color=T["grid"], heavy_color=T["heavy"])
    return ws


# ── 렌더러: Word ──────────────────────────────────────────────────────────────
def render_docx(grid, doc, *, grid_color=C_GRID, heavy_color=C_HEAVY):
    from docx.enum.table import WD_ALIGN_VERTICAL
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import Inches, Pt, RGBColor

    SZ = {THIN: (grid_color, 4), HEAVY: (heavy_color, 12)}
    HAL = {"left": WD_ALIGN_PARAGRAPH.LEFT, "right": WD_ALIGN_PARAGRAPH.RIGHT, "center": WD_ALIGN_PARAGRAPH.CENTER}

    def shade(tc, hexcolor):
        shd = OxmlElement("w:shd")
        shd.set(qn("w:val"), "clear"); shd.set(qn("w:color"), "auto"); shd.set(qn("w:fill"), hexcolor)
        tc._tc.get_or_add_tcPr().append(shd)

    def borders(tc, sides):
        tcPr = tc._tc.get_or_add_tcPr()
        tb = tcPr.find(qn("w:tcBorders"))
        if tb is None:
            tb = OxmlElement("w:tcBorders"); tcPr.append(tb)
        for s in ("top", "left", "right", "bottom"):
            el = tb.find(qn("w:" + s))
            if el is None:
                el = OxmlElement("w:" + s); tb.append(el)
            spec = sides.get(s)
            if spec:
                col, sz = SZ[spec]
                el.set(qn("w:val"), "single"); el.set(qn("w:sz"), str(sz)); el.set(qn("w:color"), col)
            else:
                el.set(qn("w:val"), "nil")

    rec = reconcile(grid)
    tbl = doc.add_table(rows=grid.nr, cols=grid.nc)
    tbl.allow_autofit = False
    tlay = OxmlElement("w:tblLayout"); tlay.set(qn("w:type"), "fixed")
    tbl._tbl.tblPr.append(tlay)
    cm = OxmlElement("w:tblCellMar")              # 여백 축소(좁은 열에서 줄바꿈 방지)
    for s, w in (("left", 55), ("right", 55), ("top", 15), ("bottom", 15)):
        e = OxmlElement("w:" + s); e.set(qn("w:w"), str(w)); e.set(qn("w:type"), "dxa"); cm.append(e)
    tbl._tbl.tblPr.append(cm)

    # Word vMerge 좌측선 버그 회피 → 세로 병합을 '행별 가로 병합'으로 분해(텍스트=첫행, 채움=전행 동일)
    merges = grid.merges + grid.soft_merges
    covered, lead_region = set(), {}
    for (r1, c1, r2, c2) in merges:
        anchor = grid.g.get((r1, c1)) or GCell()
        for r in range(r1, r2 + 1):
            lead_region[(r, c1)] = (r1, c1, r2, c2)
            if r != r1:
                grid.cell(r, c1).fill = anchor.fill
            for c in range(c1 + 1, c2 + 1):
                covered.add((r, c))
    for (r1, c1, r2, c2) in merges:
        if c1 < c2:
            for r in range(r1, r2 + 1):
                tbl.cell(r - 1, c1 - 1).merge(tbl.cell(r - 1, c2 - 1))

    for r in range(1, grid.nr + 1):
        for c in range(1, grid.nc + 1):
            if (r, c) in covered:
                continue
            cell = grid.g.get((r, c)) or GCell()
            tc = tbl.cell(r - 1, c - 1)
            tc.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
            para = tc.paragraphs[0]
            para.paragraph_format.space_after = Pt(0)
            para.paragraph_format.space_before = Pt(0)
            if cell.halign:
                para.alignment = HAL[cell.halign]
            reg = lead_region.get((r, c))
            if cell.text and (reg is None or r == reg[0]):
                run = para.add_run(cell.text)
                run.font.name = FONT
                run.font.size = Pt(cell.size)
                run.font.bold = cell.bold
                run.font.color.rgb = RGBColor.from_string(cell.color)
                rpr = run._element.get_or_add_rPr()
                rpr.get_or_add_rFonts().set(qn("w:eastAsia"), FONT)
            if cell.fill:
                shade(tc, cell.fill)
            if reg:
                br1, bc1, br2, bc2 = reg
                t, l, rt, b = _region_border(rec, br1, bc1, br2, bc2)
                borders(tc, {"top": t if r == br1 else None, "left": l,
                             "right": rt, "bottom": b if r == br2 else None})
            else:
                t, l, rt, b = _region_border(rec, r, c, r, c)
                if _is_dark(cell.fill):   # 어두운 셀: 밝은 이웃에게 공유 굵은선 양보(navy에 묻힘 방지)
                    gg = grid.g.get
                    if l == HEAVY and c > 1 and not _is_dark((gg((r, c - 1)) or GCell()).fill):
                        l = None
                    if rt == HEAVY and c < grid.nc and not _is_dark((gg((r, c + 1)) or GCell()).fill):
                        rt = None
                    if t == HEAVY and r > 1 and not _is_dark((gg((r - 1, c)) or GCell()).fill):
                        t = None
                    if b == HEAVY and r < grid.nr and not _is_dark((gg((r + 1, c)) or GCell()).fill):
                        b = None
                borders(tc, {"top": t, "left": l, "right": rt, "bottom": b})
            tc._tc.get_or_add_tcPr().append(OxmlElement("w:noWrap"))   # 숫자 한 줄 유지

    sec = doc.sections[-1]
    usable_in = (sec.page_width - sec.left_margin - sec.right_margin) / 914400
    raw = {c: grid.widths.get(c, 8.43) * 0.092 for c in range(1, grid.nc + 1)}
    scale = min(1.0, usable_in / sum(raw.values())) if raw else 1.0
    for c in range(1, grid.nc + 1):
        w = Inches(raw[c] * scale)
        for r in range(grid.nr):
            tbl.cell(r, c - 1).width = w
    return tbl


def build_docx(variants, path, theme="color"):
    from docx import Document
    from docx.shared import Inches
    T = THEME[theme]
    doc = Document()
    sec = doc.sections[0]
    sec.left_margin = sec.right_margin = Inches(0.6)
    for i, v in enumerate(variants):
        model, orient = v[0], v[1]
        render_docx(themed(layout(model, orient), theme), doc, grid_color=T["grid"], heavy_color=T["heavy"])
        if i < len(variants) - 1:
            doc.add_paragraph()
    doc.save(path)
    return path


# ── 렌더러: PowerPoint ────────────────────────────────────────────────────────
def render_pptx(grid, slide, *, left_in=0.4, top_in=0.4, max_w_in=12.5, row_h_in=0.32,
                grid_color=C_GRID, heavy_color=C_HEAVY):
    from pptx.oxml.ns import qn
    from pptx.util import Emu, Inches, Pt

    EMU = 914400
    span, covered = _merge_maps(grid.merges + grid.soft_merges)
    rec = reconcile(grid)

    col_in = {}
    for c in range(1, grid.nc + 1):
        col_in[c] = max(0.5, grid.widths.get(c, 8.43) * 0.095)
    total = sum(col_in.values())
    if total > max_w_in:
        s = max_w_in / total
        col_in = {c: w * s for c, w in col_in.items()}
        total = max_w_in

    def row_in(r):
        cell = grid.g.get((r, 1))
        sz = cell.size if cell else 10
        return 0.40 if sz >= 13 else (0.26 if sz == 9 else row_h_in)

    height = sum(row_in(r) for r in range(1, grid.nr + 1))
    gf = slide.shapes.add_table(grid.nr, grid.nc, Inches(left_in), Inches(top_in),
                                Inches(total), Inches(height))
    tbl = gf.table
    tbl.first_row = tbl.last_row = tbl.first_col = tbl.horz_banding = tbl.vert_banding = False
    tblPr = tbl._tbl.find(qn("a:tblPr"))      # 기본 표 스타일 제거 → 내 테두리만 권위(PowerPoint 덧칠 방지)
    if tblPr is not None:
        sid = tblPr.find(qn("a:tableStyleId"))
        if sid is None:
            sid = tblPr.makeelement(qn("a:tableStyleId"), {})
            tblPr.append(sid)
        sid.text = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"   # No Style, No Grid

    for c in range(1, grid.nc + 1):
        tbl.columns[c - 1].width = Emu(int(col_in[c] * EMU))
    for r in range(1, grid.nr + 1):
        tbl.rows[r - 1].height = Emu(int(row_in(r) * EMU))

    for (r1, c1), (r2, c2) in span.items():
        tbl.cell(r1 - 1, c1 - 1).merge(tbl.cell(r2 - 1, c2 - 1))

    def set_borders(tc, sides):
        tcPr = tc._tc.get_or_add_tcPr()
        for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
            for e in tcPr.findall(qn(tag)):
                tcPr.remove(e)
        order = (("a:lnL", sides[1]), ("a:lnR", sides[2]), ("a:lnT", sides[0]), ("a:lnB", sides[3]))
        for i, (tag, spec) in enumerate(order):
            ln = tcPr.makeelement(qn(tag), {"cap": "flat"})
            if spec:
                col, emu = (heavy_color, 19050) if spec == HEAVY else (grid_color, 6350)
                ln.set("w", str(emu))
                fill = tcPr.makeelement(qn("a:solidFill"), {})
                clr = tcPr.makeelement(qn("a:srgbClr"), {"val": col})
                fill.append(clr); ln.append(fill)
            else:
                ln.set("w", "3175")
                ln.append(tcPr.makeelement(qn("a:noFill"), {}))
            tcPr.insert(i, ln)

    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.dml.color import RGBColor
    HAL = {"left": PP_ALIGN.LEFT, "right": PP_ALIGN.RIGHT, "center": PP_ALIGN.CENTER}
    for r in range(1, grid.nr + 1):
        for c in range(1, grid.nc + 1):
            if (r, c) in covered:
                continue
            cell = grid.g.get((r, c)) or GCell()
            tc = tbl.cell(r - 1, c - 1)
            tc.margin_left = tc.margin_right = Inches(0.04)
            tc.margin_top = tc.margin_bottom = Inches(0.01)
            tc.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = tc.text_frame
            tf.word_wrap = bool(cell.wrap)
            p = tf.paragraphs[0]
            if cell.halign:
                p.alignment = HAL[cell.halign]
            run = p.add_run()
            run.text = cell.text
            run.font.size = Pt(cell.size)
            run.font.bold = cell.bold
            run.font.name = FONT
            run.font.color.rgb = RGBColor.from_string(cell.color)
            if cell.fill:
                tc.fill.solid(); tc.fill.fore_color.rgb = RGBColor.from_string(cell.fill)
            else:
                tc.fill.background()
            r2, c2 = span.get((r, c), (r, c))
            set_borders(tc, _region_border(rec, r, c, r2, c2))
    return gf


def build_pptx(variants, path, theme="color"):
    from pptx import Presentation
    from pptx.util import Inches
    T = THEME[theme]
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for v in variants:
        model, orient = v[0], v[1]
        slide = prs.slides.add_slide(blank)
        render_pptx(themed(layout(model, orient), theme), slide, grid_color=T["grid"], heavy_color=T["heavy"])
    prs.save(path)
    return path


# ── 내장 샘플 ─────────────────────────────────────────────────────────────────
SAMPLE = {
    "title": "2024년 사업부별 실적",
    "unit": "단위: 억원 / %",
    "corner_label": "사업부",
    "attributes": ["매출", "영업이익", "점유율"],
    "items": [
        {"label": "A사업부", "summary": [1200, 180, "21%"], "details": [
            {"label": "Q1", "values": [280, 38, "20%"]},
            {"label": "Q2", "values": [300, 42, "21%"]},
            {"label": "Q3", "values": [300, 48, "22%"]},
            {"label": "Q4", "values": [320, 52, "21%"]},
        ]},
        {"label": "B사업부", "summary": [800, 90, "12%"], "details": [
            {"label": "Q1", "values": [180, 18, "11%"]},
            {"label": "Q2", "values": [200, 22, "12%"]},
            {"label": "Q3", "values": [210, 24, "12%"]},
            {"label": "Q4", "values": [210, 26, "13%"]},
        ]},
        {"label": "C사업부", "summary": [2100, 520, "33%"], "details": [
            {"label": "Q1", "values": [480, 110, "31%"]},
            {"label": "Q2", "values": [510, 125, "32%"]},
            {"label": "Q3", "values": [540, 135, "34%"]},
            {"label": "Q4", "values": [570, 150, "35%"]},
        ]},
    ],
}

SAMPLE_2D = {
    "title": "사업부 × 연도 매출",
    "unit": "단위: 억원",
    "corner_label": "사업부 ＼ 연도",
    "row_groups": [
        {"label": "A사업부", "details": ["국내", "해외"]},
        {"label": "B사업부", "details": ["국내", "해외"]},
    ],
    "col_groups": [
        {"label": "2023", "details": ["상반기", "하반기"]},
        {"label": "2024", "details": ["상반기", "하반기"]},
    ],
    "data": {
        "A사업부": {
            "국내": {"2023": {"상반기": 700, "하반기": 750}, "2024": {"상반기": 780, "하반기": 820}},
            "해외": {"2023": {"상반기": 450, "하반기": 500}, "2024": {"상반기": 470, "하반기": 530}},
        },
        "B사업부": {
            "국내": {"2023": {"상반기": 300, "하반기": 320}, "2024": {"상반기": 310, "하반기": 330}},
            "해외": {"2023": {"상반기": 250, "하반기": 260}, "2024": {"상반기": 270, "하반기": 280}},
        },
    },
}

_VARIANTS = [(SAMPLE, "row"), (SAMPLE, "column"), (SAMPLE_2D, "both")]


# ── CLI ───────────────────────────────────────────────────────────────────────
def load_data(path):
    """데이터 로드 — .yaml/.yml 은 YAML, 그 외 JSON."""
    text = open(path, encoding="utf-8").read()
    if path.lower().endswith((".yaml", ".yml")):
        import yaml
        return yaml.safe_load(text)
    return json.loads(text)


def generate(data, path, orientation="row", theme="color"):
    """단일 표 1개 생성 — 출력 확장자(.xlsx/.docx/.pptx)로 포맷 자동 판별."""
    ext = path.rsplit(".", 1)[-1].lower()
    if ext == "xlsx":
        build_xlsx(data, path, orientation, theme)
    elif ext == "docx":
        build_docx([(data, orientation)], path, theme)
    elif ext == "pptx":
        build_pptx([(data, orientation)], path, theme)
    else:
        raise ValueError("출력 확장자는 .xlsx / .docx / .pptx 중 하나여야 합니다")
    return path


_USAGE = ("drilldown-table CLI\n"
          "  py drilldown_table.py DATA OUT [orient] [theme]\n"
          "    DATA   : data file (.json or .yaml)\n"
          "    OUT    : output (.xlsx | .docx | .pptx -- format from extension)\n"
          "    orient : row | column | both     (default row)\n"
          "    theme  : color | grey | mono      (default color)\n"
          "  ex) py drilldown_table.py data.yaml out.pptx both grey\n"
          "  no args -> demo: drilldown_sample.{xlsx,docx,pptx}")


if __name__ == "__main__":
    a = sys.argv[1:]
    if a and a[0] in ("-h", "--help"):
        print(_USAGE)
    elif len(a) >= 2:
        orient = a[2] if len(a) > 2 else "row"
        theme = a[3] if len(a) > 3 else "color"
        generate(load_data(a[0]), a[1], orient, theme)
        print(f"saved {a[1]}  (orient={orient}, theme={theme})")
    else:
        base = "drilldown_sample"
        build_workbook(_VARIANTS, base + ".xlsx")
        build_docx(_VARIANTS, base + ".docx")
        build_pptx(_VARIANTS, base + ".pptx")
        print(f"saved {base}.xlsx / .docx / .pptx  (데모: 행/열/행+열 × color)")
